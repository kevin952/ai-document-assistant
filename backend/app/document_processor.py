import os
from typing import List, Optional
from fastapi import UploadFile
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import ChatOpenAI
from langchain_core.documents import Document
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from dotenv import load_dotenv
import docx
from PyPDF2 import PdfReader
from pptx import Presentation
import markdown
import tempfile
import logging

# Configure logging to only show errors
logging.basicConfig(level=logging.ERROR)
logger = logging.getLogger(__name__)

load_dotenv()

class DocumentProcessor:
    def __init__(self):
        logger.info("Initializing DocumentProcessor...")
        self.embeddings = OpenAIEmbeddings(
            openai_api_key=os.getenv("OPENAI_API_KEY")
        )
        logger.info("Initialized OpenAI embeddings")
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        logger.info("Initialized text splitter")
        
        self.vectorstore_dir = "data/vectorstore"
        os.makedirs(self.vectorstore_dir, exist_ok=True)
        logger.info(f"Created vectorstore directory: {self.vectorstore_dir}")
        
        self.vectorstore = Chroma(
            persist_directory=self.vectorstore_dir,
            embedding_function=self.embeddings,
            collection_name="company_documents"
        )
        logger.info("Initialized Chroma vectorstore")
        
        # Track processed files
        self.processed_files = set()
        self._load_processed_files()

    def _load_processed_files(self):
        """Load list of already processed files"""
        try:
            if os.path.exists(os.path.join(self.vectorstore_dir, "processed_files.txt")):
                with open(os.path.join(self.vectorstore_dir, "processed_files.txt"), "r") as f:
                    self.processed_files = set(f.read().splitlines())
        except Exception as e:
            logger.error(f"Error loading processed files list: {str(e)}")

    def _save_processed_files(self):
        """Save list of processed files"""
        try:
            with open(os.path.join(self.vectorstore_dir, "processed_files.txt"), "w") as f:
                f.write("\n".join(self.processed_files))
        except Exception as e:
            logger.error(f"Error saving processed files list: {str(e)}")

    async def process_document(self, file: UploadFile):
        try:
            logger.info(f"Processing uploaded file: {file.filename}")
            # Save the uploaded file temporarily
            file_path = f"data/{file.filename}"
            with open(file_path, "wb") as buffer:
                content = await file.read()
                buffer.write(content)
            
            # Extract text based on file type
            text = self._extract_text(file_path)
            logger.info(f"Extracted {len(text)} characters from {file.filename}")
            
            # Create a document
            doc = Document(page_content=text, metadata={"source": file.filename})
            
            # Split the document
            splits = self.text_splitter.split_documents([doc])
            logger.info(f"Split {file.filename} into {len(splits)} chunks")
            
            # Add to vectorstore
            self.vectorstore.add_documents(splits)
            logger.info(f"Added {file.filename} to vectorstore")
            
            # Clean up the temporary file
            os.remove(file_path)
            logger.info(f"Cleaned up temporary file: {file_path}")
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            raise

    def _extract_text(self, file_path: str) -> str:
        """Extract text from different file formats"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.docx':
            doc = docx.Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        
        elif file_extension == '.pdf':
            reader = PdfReader(file_path)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n'
            return text
        
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
            return text
        
        elif file_extension == '.md':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")

    def process_all_documents(self):
        """Process all documents in the data directory"""
        logger.info("Starting to process all documents...")
        data_dir = os.path.join(os.path.dirname(__file__), "data")
        supported_extensions = {'.txt', '.md', '.docx', '.pdf', '.pptx'}
        
        if not os.path.exists(data_dir):
            logger.error(f"Data directory not found: {data_dir}")
            return
            
        files = [f for f in os.listdir(data_dir) if any(f.endswith(ext) for ext in supported_extensions)]
        logger.info(f"Found {len(files)} documents to process")
        
        # Filter out already processed files
        new_files = [f for f in files if f not in self.processed_files]
        logger.info(f"Found {len(new_files)} new documents to process")
        
        for filename in new_files:
            file_path = os.path.join(data_dir, filename)
            try:
                logger.info(f"Processing document: {filename}")
                # Extract text
                text = self._extract_text(file_path)
                logger.info(f"Extracted {len(text)} characters from {filename}")
                
                # Create a document
                doc = Document(page_content=text, metadata={"source": filename})
                
                # Split the document
                splits = self.text_splitter.split_documents([doc])
                logger.info(f"Split {filename} into {len(splits)} chunks")
                
                # Add to vectorstore
                self.vectorstore.add_documents(splits)
                logger.info(f"Added {filename} to vectorstore")
                
                # Mark file as processed
                self.processed_files.add(filename)
                
            except Exception as e:
                logger.error(f"Error processing {filename}: {str(e)}")
                continue
        
        # Save the updated list of processed files
        self._save_processed_files()
        logger.info("Finished processing all documents") 